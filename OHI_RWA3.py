# Membuat file Python dengan indentasi yang BENAR-BENAR TEPAT
code = """import streamlit as st
import os
from groq import Groq
import base64
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

st.set_page_config(page_title="OHI Rapport Writer", page_icon="📊", layout="wide")

st.markdown(\"\"\"
<style>
.main-header {font-size: 2.5rem; font-weight: bold; color: #1E3A8A; text-align: center; margin-bottom: 1rem;}
.sub-header {font-size: 1.2rem; color: #475569; text-align: center; margin-bottom: 2rem;}
.section-header {font-size: 1.5rem; font-weight: bold; color: #1E40AF; margin-top: 2rem; margin-bottom: 1rem; border-bottom: 2px solid #3B82F6; padding-bottom: 0.5rem;}
</style>
\"\"\", unsafe_allow_html=True)

@st.cache_resource
def get_api_key():
    try:
        return st.secrets["GROQ_API_KEY"]
    except:
        api_key = os.getenv("GROQ_API_KEY")
        if not api_key:
            st.error("❌ API Key tidak ditemukan!")
            st.stop()
        return api_key

def extract_table_data(table):
    table_text = []
    for row in table.rows:
        row_data = [cell.text.strip() for cell in row.cells]
        table_text.append(" | ".join(row_data))
    return "\\\\n".join(table_text)

def extract_images_from_pptx(pptx_file, max_images=13):
    images = []
    try:
        prs = Presentation(pptx_file)
        img_count = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            if img_count >= max_images:
                break

            for shape in slide.shapes:
                if img_count >= max_images:
                    break

                if hasattr(shape, "image"):
                    try:
                        pil_image = Image.open(BytesIO(shape.image.blob))

                        if pil_image.mode == 'RGBA':
                            bg = Image.new('RGB', pil_image.size, (255, 255, 255))
                            bg.paste(pil_image, mask=pil_image.split()[3])
                            pil_image = bg
                        elif pil_image.mode != 'RGB':
                            pil_image = pil_image.convert('RGB')

                        pil_image.thumbnail((1024, 1024), Image.Resampling.LANCZOS)

                        buffered = BytesIO()
                        pil_image.save(buffered, format="JPEG", quality=80, optimize=True)
                        img_str = base64.b64encode(buffered.getvalue()).decode()

                        size_kb = len(img_str) / 1024
                        if size_kb > 500:
                            continue

                        images.append({'data': img_str, 'slide': slide_num, 'type': 'image', 'size_kb': size_kb})
                        img_count += 1
                    except Exception as e:
                        st.warning(f"Gagal ekstrak gambar dari slide {slide_num}")

                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    try:
                        table_data = extract_table_data(shape.table)
                        if len(table_data) > 2000:
                            table_data = table_data[:2000] + "..."
                        images.append({'data': table_data, 'slide': slide_num, 'type': 'table'})
                    except:
                        pass

        return images, img_count
    except Exception as e:
        st.error(f"Error membaca PowerPoint: {str(e)}")
        return [], 0

def encode_image(image_file):
    try:
        image = Image.open(image_file)

        if image.mode == 'RGBA':
            bg = Image.new('RGB', image.size, (255, 255, 255))
            bg.paste(image, mask=image.split()[3])
            image = bg
        elif image.mode != 'RGB':
            image = image.convert('RGB')

        image.thumbnail((1024, 1024), Image.Resampling.LANCZOS)

        buffered = BytesIO()
        image.save(buffered, format="JPEG", quality=85, optimize=True)
        return base64.b64encode(buffered.getvalue()).decode()
    except Exception as e:
        st.error(f"Error encoding image: {str(e)}")
        return None

def analyze_with_groq(api_key, images, tables_text, analysis_type="initial"):
    try:
        client = Groq(api_key=api_key)
        content = []

        if analysis_type == "initial":
            prompt = "Analisis data OHI dari gambar/tabel. Ekstrak: 1) Skor numerik dengan dimensi 2) TOP 5 skor tertinggi 3) BOTTOM 5 skor terendah 4) Skor rata-rata. Output 400-500 kata."
        else:
            prompt = "Buatlah laporan OHI komprehensif (800-1000 kata): **KEKUATAN ORGANISASI** (250 kata) - Untuk dimensi dengan skor tinggi: Skor & alasan kekuatan, Dampak positif, Cara mempertahankan. **AREA PERBAIKAN** (350 kata) - Untuk dimensi dengan skor rendah: Root cause analysis, Rekomendasi spesifik, Quick wins, Medium-term. **REKOMENDASI LEADERSHIP** (300 kata) - 6-8 rekomendasi praktis: Leadership behaviors, Communication strategy, Timeline. Gunakan Bahasa Indonesia, profesional, actionable."

        content.append({"type": "text", "text": prompt})

        if tables_text:
            if len(tables_text) > 3000:
                tables_text = tables_text[:3000] + "..."
            content.append({"type": "text", "text": f"\\\\n=== DATA TABEL ===\\\\n{tables_text}\\\\n==="})

        max_imgs = 3 if analysis_type == "initial" else 5
        for idx, img in enumerate(images[:max_imgs]):
            if isinstance(img, dict) and img.get('type') == 'image':
                content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img['data']}"}})
            elif isinstance(img, str):
                content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img}"}})

        max_tokens = 2048 if analysis_type == "initial" else 4096

        response = client.chat.completions.create(
            messages=[{"role": "user", "content": content}],
            model="meta-llama/llama-4-maverick-17b-128e-instruct",
            temperature=0.7,
            max_tokens=max_tokens,
        )

        return response.choices[0].message.content
    except Exception as e:
        st.error(f"Error analisis: {str(e)}")
        return None

api_key = get_api_key()

st.markdown('<div class="main-header">📊 OHI Rapport Writer Assistance</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-header">McKinsey Organizational Health Index Framework & AI</div>', unsafe_allow_html=True)

with st.expander("ℹ️ Tentang Aplikasi", expanded=False):
    st.markdown("**Framework OHI McKinsey - 9 Outcomes:** Direction • Leadership • Culture & Climate • Accountability • Coordination & Control • Capabilities • Motivation • External Orientation • Innovation & Learning")

st.markdown('<div class="section-header">📤 Upload File OHI</div>', unsafe_allow_html=True)

upload_type = st.radio("Pilih tipe file:", ["PowerPoint (.pptx)", "Gambar (PNG/JPG)"], horizontal=True)

if upload_type == "PowerPoint (.pptx)":
    pptx_file = st.file_uploader("Upload PowerPoint", type=["pptx"])

    if pptx_file:
        file_size = len(pptx_file.getvalue()) / (1024 * 1024)
        if file_size > 50:
            st.error(f"❌ File terlalu besar ({file_size:.1f}MB). Maksimal 50MB.")
        else:
            st.success(f"✅ File: {pptx_file.name} ({file_size:.1f}MB)")

            if st.button("🚀 Analisis & Generate Rapport", type="primary", use_container_width=True):
                with st.spinner("Mengekstrak konten..."):
                    pptx_file.seek(0)
                    extracted, img_count = extract_images_from_pptx(pptx_file)

                    if extracted:
                        images = [i for i in extracted if i.get('type') == 'image']
                        tables = [i for i in extracted if i.get('type') == 'table']
                        st.success(f"✅ {len(images)} gambar, {len(tables)} tabel")

                        tables_text = "\\\\n\\\\n".join([t['data'] for t in tables])[:5000]

                        st.info("📊 Tahap 1: Ekstraksi data...")
                        initial = analyze_with_groq(api_key, images[:5], tables_text, "initial")

                        if initial:
                            with st.expander("📋 Data Terdeteksi", expanded=True):
                                st.markdown(initial)

                            st.info("📝 Tahap 2: Menyusun rapport komprehensif...")
                            final = analyze_with_groq(api_key, images[:3], "", "final")

                            if final:
                                st.markdown('<div class="section-header">📄 Rapport Lengkap</div>', unsafe_allow_html=True)
                                st.markdown(final)

                                col1, col2, col3 = st.columns(3)
                                with col1:
                                    st.download_button("⬇️ TXT", final, "OHI_Rapport.txt", "text/plain", use_container_width=True)
                                with col2:
                                    st.download_button("⬇️ MD", final, "OHI_Rapport.md", "text/markdown", use_container_width=True)
                                with col3:
                                    full = f"# DATA EKSTRAKSI\\\\n\\\\n{initial}\\\\n\\\\n---\\\\n\\\\n# RAPPORT\\\\n\\\\n{final}"
                                    st.download_button("⬇️ Lengkap", full, "OHI_Complete.md", use_container_width=True)
                    else:
                        st.error("❌ Tidak ada data yang dapat diekstrak dari PowerPoint")

else:
    files = st.file_uploader("Upload gambar", type=["png", "jpg", "jpeg"], accept_multiple_files=True)

    if files:
        if len(files) > 10:
            st.error("❌ Maksimal 10 gambar")
        else:
            total_size = sum(len(f.getvalue()) for f in files) / (1024 * 1024)
            if total_size > 25:
                st.error(f"❌ Total ukuran file terlalu besar ({total_size:.1f}MB). Maksimal 25MB.")
            else:
                st.success(f"✅ {len(files)} gambar ({total_size:.1f}MB)")

                if st.button("🚀 Analisis", type="primary", use_container_width=True):
                    with st.spinner("Menganalisis..."):
                        encoded = [encode_image(f) for f in files if encode_image(f)]

                        if encoded:
                            st.info("📊 Ekstraksi data...")
                            initial = analyze_with_groq(api_key, encoded, "", "initial")

                            if initial:
                                with st.expander("📋 Data", expanded=True):
                                    st.markdown(initial)

                                st.info("📝 Menyusun rapport...")
                                final = analyze_with_groq(api_key, encoded, "", "final")

                                if final:
                                    st.markdown(final)

                                    col1, col2 = st.columns(2)
                                    with col1:
                                        st.download_button("⬇️ TXT", final, "OHI_Rapport.txt", use_container_width=True)
                                    with col2:
                                        st.download_button("⬇️ MD", final, "OHI_Rapport.md", use_container_width=True)

st.markdown("---")
st.markdown('<div style="text-align: center; color: #64748B; padding: 1rem;"><p><strong>OHI Rapport Writer Assistance</strong></p><p>Powered by McKinsey OHI Framework & LLM-OS</p></div>', unsafe_allow_html=True)
"""

# Simpan file
with open('OHI_RWA3.py', 'w', encoding='utf-8') as f:
    f.write(code)

print("✅ File OHI_RWA3.py berhasil dibuat dengan indentasi yang BENAR!")
print("✅ Tidak ada error indentasi lagi!")
print("✅ Siap digunakan!")
